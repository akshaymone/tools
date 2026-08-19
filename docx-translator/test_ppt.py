import sys
from pathlib import Path
from pptx import Presentation

# Create a test presentation
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
title.text = "Hello 안녕"

# Add an image (mock)
# we need a dummy image
with open('dummy.png', 'wb') as f:
    f.write(b'\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01\x08\x06\x00\x00\x00\x1f\x15\xc4\x89\x00\x00\x00\nIDATx\x9cc\x00\x01\x00\x00\x05\x00\x01\x0d\n-\xb4\x00\x00\x00\x00IEND\xaeB`\x82')
slide.shapes.add_picture('dummy.png', 0, 0, 100, 100)

# Add notes slide
notes_slide = slide.notes_slide
text_frame = notes_slide.notes_text_frame
text_frame.text = "This is a note."

prs.save("test_input.pptx")
print("Saved test_input.pptx")
