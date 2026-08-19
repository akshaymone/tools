from pptx import Presentation
prs = Presentation()
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
notes_slide = slide.notes_slide
prs.save('test_pptx_notes.pptx')
