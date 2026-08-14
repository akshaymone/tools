"""
translator/pptx_handler.py

Walks every slide in a Presentation and translates:
  - Text frames (titles, body, text boxes)
  - Tables (cell by cell)
  - Group shapes (recursively unwrapped)
  - Speaker notes
  - Embedded picture shapes → OCR significant text → append to speaker notes
    (images are NEVER modified; translation is placed in the notes pane)

Formatting contract:
  Translation is done at the paragraph level.  The translated string is
  written into the FIRST non-empty run of each paragraph; all subsequent
  runs in that paragraph are cleared.  This preserves the paragraph's
  font, size, bold, italic, colour, and alignment properties.

  For tables, the same logic applies per cell.

  For images: only text blocks >= min_height pixels tall are translated
  (skips tiny axis labels, watermarks, etc.).  All translations from all
  images on a slide are collected and appended to the slide's speaker notes
  under a clearly labelled section.
"""

import logging
from pathlib import Path
from typing import Optional

log = logging.getLogger(__name__)

# MSO shape type constants (python-pptx / OOXML)
_MSO_GROUP   = 6   # MsoShapeType.GROUP
_MSO_PICTURE = 13  # MsoShapeType.PICTURE

# Minimum OCR word-height (pixels, after upscaling) to include in notes.
# 18 px at 2× upscale ≈ 9 px original ≈ ~14 pt — "body text" size.
# Increase this to skip smaller text.
_MIN_TEXT_HEIGHT_PX = 18


class PPTXHandler:
    def __init__(
        self,
        engine,
        skip_images: bool = False,
        confidence: int = 60,
        ocr_lang: str = "kor",
        dry_run: bool = False,
        min_text_height: int = _MIN_TEXT_HEIGHT_PX,
    ) -> None:
        self.engine = engine
        self.dry_run = dry_run
        self.min_text_height = min_text_height

        if skip_images:
            self.image_handler: Optional[object] = None
        else:
            from translator.image_handler import ImageHandler
            self.image_handler = ImageHandler(
                engine=engine,
                confidence=confidence,
                ocr_lang=ocr_lang,
            )

    # ------------------------------------------------------------------
    # Public entry point
    # ------------------------------------------------------------------

    def translate_file(self, input_path: Path, output_path: Path) -> None:
        from pptx import Presentation

        prs = Presentation(str(input_path))
        n = len(prs.slides)

        for idx, slide in enumerate(prs.slides, 1):
            log.info(f"   Slide {idx}/{n}")
            self._process_slide(slide)

        if not self.dry_run:
            output_path.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(output_path))

    # ------------------------------------------------------------------
    # Slide-level processing
    # ------------------------------------------------------------------

    def _process_slide(self, slide) -> None:
        # Collect image OCR translations first (before touching notes)
        image_translations: list[str] = []

        for shape in slide.shapes:
            self._process_shape(shape, image_translations)

        # Translate existing speaker notes text
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            self._translate_text_frame(notes_tf, label="[notes]")

        # Append image translations to speaker notes
        if image_translations and not self.dry_run:
            self._append_image_notes(slide, image_translations)
        elif image_translations and self.dry_run:
            log.info(f"      [DRY] Would append {len(image_translations)} image translation(s) to notes.")

    def _append_image_notes(self, slide, translations: list[str]) -> None:
        """Append translated image text to the slide's speaker notes pane."""
        from pptx.util import Pt
        from lxml import etree

        # Ensure notes slide exists
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame

        # Add a blank separator paragraph
        sep_para = tf.add_paragraph()
        sep_para.text = ""

        # Header paragraph
        header_para = tf.add_paragraph()
        header_para.text = "── Image Text (auto-translated) ──"
        if header_para.runs:
            header_para.runs[0].font.bold = True
            header_para.runs[0].font.size = Pt(9)

        # One paragraph per translated block
        for i, text in enumerate(translations, 1):
            p = tf.add_paragraph()
            p.text = f"{i}. {text}"
            if p.runs:
                p.runs[0].font.size = Pt(9)

        log.debug(f"      Appended {len(translations)} image translation(s) to notes.")

    # ------------------------------------------------------------------
    # Shape dispatch
    # ------------------------------------------------------------------

    def _process_shape(self, shape, image_translations: list) -> None:
        shape_type = shape.shape_type

        # Recursively unwrap group shapes
        if shape_type == _MSO_GROUP:
            for child in shape.shapes:
                self._process_shape(child, image_translations)
            return

        # Text frames (text boxes, placeholders, auto-shapes with text)
        if shape.has_text_frame:
            self._translate_text_frame(shape.text_frame, label=shape.name)

        # Tables
        if shape.has_table:
            self._translate_table(shape.table)

        # Embedded pictures → OCR → speaker notes (image NOT modified)
        if self.image_handler is not None and shape_type == _MSO_PICTURE:
            self._process_picture(shape, image_translations)

    # ------------------------------------------------------------------
    # Text frame / paragraph translation
    # ------------------------------------------------------------------

    def _translate_text_frame(self, text_frame, label: str = "") -> None:
        for para in text_frame.paragraphs:
            self._translate_paragraph(para, label=label)

    def _translate_paragraph(self, para, label: str = "") -> None:
        full_text = para.text.strip()
        if not full_text:
            return

        translated = self.engine.translate(full_text)
        if not translated or translated.strip() == full_text:
            return

        if self.dry_run:
            log.info(f"      [DRY] {label}: {full_text[:55]!r}")
            log.info(f"              → {translated[:55]!r}")
            return

        # Collect all runs that have content
        active_runs = [r for r in para.runs if r.text.strip()]
        if not active_runs:
            return

        # Write translated text into first run, blank out the rest
        active_runs[0].text = translated
        for run in active_runs[1:]:
            run.text = ""

        log.debug(f"      {label}: {full_text[:40]!r} → {translated[:40]!r}")

    def _translate_table(self, table) -> None:
        for row in table.rows:
            for cell in row.cells:
                self._translate_text_frame(cell.text_frame, label="[table]")

    # ------------------------------------------------------------------
    # Image processing → speaker notes
    # ------------------------------------------------------------------

    def _process_picture(self, shape, image_translations: list) -> None:
        """
        OCR the image for significant text, translate, collect into
        image_translations list.  The image itself is NOT modified.
        """
        try:
            image = shape.image
        except Exception as exc:
            log.warning(f"      Could not access image for shape {shape.name!r}: {exc}")
            return

        size_kb = len(image.blob) // 1024
        log.debug(f"      OCR → shape {shape.name!r} ({size_kb} KB, {image.content_type})")

        if self.dry_run:
            log.info(f"      [DRY] Would OCR image shape {shape.name!r} for speaker notes.")
            return

        try:
            texts = self.image_handler.extract_text_for_notes(
                image.blob,
                min_height=self.min_text_height,
            )
        except Exception as exc:
            log.warning(f"      Image OCR failed [{shape.name}]: {exc}")
            return

        if texts:
            log.debug(f"      Found {len(texts)} translatable block(s) in {shape.name!r}.")
            image_translations.extend(texts)
        else:
            log.debug(f"      No considerable Korean text found in {shape.name!r}.")
