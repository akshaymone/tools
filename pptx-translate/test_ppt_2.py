import sys
from pptx import Presentation

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
title.text = "Hello 안녕"
prs.save("test_namespaces.pptx")
