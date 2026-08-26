import sys
from pptx import Presentation

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
# Initialize notes slide without text
_ = slide.notes_slide
prs.save("test_empty_notes.pptx")
print("Saved")
